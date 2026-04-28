# -*- coding: utf-8 -*-
"""Build thesis poster following the original ISIS poster template structure."""

from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt


ROOT = Path(__file__).resolve().parent
TEMPLATE = Path(r"C:\Users\VivoBook\Downloads\FormatoPoster_ISIS_v2.potx")
ASSET_DIR = ROOT / "images" / "Poster"

OUT_PPTX = ROOT / "Poster_Tesis_HFL_v7.pptx"
OUT_PREVIEW = ROOT / "Poster_Tesis_HFL_v7_preview.png"

HEADER_ORANGE = ASSET_DIR / "template_header_orange.png"
BAR_BLUE_FULL = ASSET_DIR / "template_bar_blue_full.png"
BAR_ORANGE_HALF = ASSET_DIR / "template_bar_orange_half.png"
BAR_BLUE_HALF = ASSET_DIR / "template_bar_blue_half.png"

ACC_LOSS_IMG = ROOT / "images" / "Resultados" / "hfl_v7_accuracy_loss_matrix.png"
WEIGHTS_IMG = ROOT / "images" / "Resultados" / "hfl_v7_weight_magnitude_matrix.png"

EMU_PER_IN = 914400
SLIDE_W = 27.56
SLIDE_H = 39.38


class C:
    ORANGE = RGBColor(239, 105, 34)
    BLUE = RGBColor(37, 91, 154)
    INK = RGBColor(30, 34, 38)
    MUTED = RGBColor(91, 99, 106)
    DASH = RGBColor(120, 120, 120)
    PAPER = RGBColor(250, 248, 244)
    WHITE = RGBColor(255, 255, 255)
    PEACH = RGBColor(241, 218, 206)
    GREEN = RGBColor(0, 105, 55)


def inch(v: float) -> int:
    return int(v * EMU_PER_IN)


def rgb_hex(rgb: RGBColor) -> str:
    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"


def extract_template_assets() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    assets = {
        HEADER_ORANGE: "ppt/media/image2.png",
        BAR_BLUE_FULL: "ppt/media/image1.png",
        BAR_ORANGE_HALF: "ppt/media/image3.png",
        BAR_BLUE_HALF: "ppt/media/image4.png",
    }
    with ZipFile(TEMPLATE) as zf:
        for out, member in assets.items():
            if not out.exists():
                out.write_bytes(zf.read(member))


def set_text(
    shape,
    text: str,
    size: float,
    color: RGBColor = C.INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Segoe UI",
    valign=MSO_ANCHOR.TOP,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.02)
    tf.margin_right = inch(0.02)
    tf.margin_top = inch(0.01)
    tf.margin_bottom = inch(0.01)
    tf.vertical_anchor = valign
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(0)


def add_text(slide, x, y, w, h, text, size, color=C.INK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    set_text(tb, text, size, color, bold, align, valign=valign)
    return tb


def add_rect(slide, x, y, w, h, fill=C.WHITE, line=C.DASH, dashed=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.0)
    if dashed:
        shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shp


def add_section(slide, x, y, w, h, title, color="orange"):
    bar_h = 0.66
    bar = BAR_ORANGE_HALF if color == "orange" else BAR_BLUE_HALF
    if w > 20:
        bar = BAR_BLUE_FULL if color == "blue" else BAR_ORANGE_HALF
    slide.shapes.add_picture(str(bar), inch(x), inch(y), width=inch(w), height=inch(bar_h))
    add_text(slide, x + 0.28, y + 0.14, w - 0.56, 0.32, title, 15.2, C.WHITE, True)
    add_rect(slide, x, y + bar_h, w, h - bar_h, C.WHITE, C.DASH, dashed=True)
    return x + 0.36, y + 0.98, w - 0.72, h - 1.22


def add_bullets(slide, x, y, w, h, lines, size=8.8):
    add_text(slide, x, y, w, h, "\n".join(f"• {line}" for line in lines), size)


def add_image_fit(slide, image_path: Path, x, y, w, h):
    img = Image.open(image_path)
    iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        width = w
        height = w / img_ratio
        xx = x
        yy = y + (h - height) / 2
    else:
        height = h
        width = h * img_ratio
        xx = x + (w - width) / 2
        yy = y
    slide.shapes.add_picture(str(image_path), inch(xx), inch(yy), width=inch(width), height=inch(height))


def mini_architecture(slide, x, y, w, h):
    def box(px, py, pw, ph, label, fill):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(px), inch(py), inch(pw), inch(ph))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = C.DASH
        add_text(slide, px + 0.08, py + 0.12, pw - 0.16, ph - 0.12, label, 10.2, C.INK, True, PP_ALIGN.CENTER)

    pc = (x + w * 0.33, y + 0.10, w * 0.34, 0.90)
    leader = (x + w * 0.06, y + 1.50, w * 0.38, 0.90)
    peer = (x + w * 0.56, y + 1.50, w * 0.38, 0.90)
    e1 = (x + w * 0.06, y + 3.00, w * 0.38, 0.90)
    e2 = (x + w * 0.56, y + 3.00, w * 0.38, 0.90)
    box(*pc, "PC\nCloud + FedAvg", RGBColor(235, 242, 252))
    box(*leader, "RPi\nFOG leader", RGBColor(238, 248, 239))
    box(*peer, "RPi\nFOG peer", RGBColor(238, 248, 239))
    box(*e1, "ESP32-S3\nEdge A", RGBColor(253, 239, 229))
    box(*e2, "ESP32-S3\nEdge B", RGBColor(253, 239, 229))

    # Use block arrows to avoid renderer-specific connector behavior.
    for px, py, color in [
        (x + w * 0.49, y + 1.05, C.ORANGE),
        (x + w * 0.24, y + 2.45, C.ORANGE),
        (x + w * 0.74, y + 2.45, C.ORANGE),
    ]:
        shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, inch(px), inch(py), inch(0.32), inch(0.48))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.color.rgb = color

    add_text(slide, x, y + 4.10, w, 0.42, "Top-down: modelo global. Bottom-up: features, pesos locales y FedAvg.", 9.0, C.MUTED, False, PP_ALIGN.CENTER)


def small_results_table(slide, x, y, w, h):
    rows = [
        ("Variante", "Acc.", "Loss", "Ronda"),
        ("RN + ASCON", "94.63%", "0.1415", "61.25 s"),
        ("RN sin ASCON", "92.98%", "0.1581", "51.01 s"),
        ("CNN-1D + FOG", "97.50%", "0.0910", "58.61 s"),
    ]
    col = [0.43, 0.18, 0.18, 0.21]
    rh = h / len(rows)
    for r, row in enumerate(rows):
        xx = x
        for c, value in enumerate(row):
            cw = w * col[c]
            fill = C.BLUE if r == 0 else (RGBColor(238, 242, 247) if r % 2 else C.WHITE)
            text_color = C.WHITE if r == 0 else C.INK
            add_rect(slide, xx, y + r * rh, cw, rh, fill, RGBColor(190, 190, 190), dashed=False)
            add_text(slide, xx + 0.06, y + r * rh + 0.07, cw - 0.12, rh - 0.06, value, 9.3, text_color, r == 0)
            xx += cw


def accuracy_bars(slide, x, y, w, h):
    data = [
        ("RN + ASCON", 0.9463, C.BLUE),
        ("RN sin ASCON", 0.9298, RGBColor(45, 155, 128)),
        ("CNN-1D + FOG", 0.9750, C.ORANGE),
    ]
    add_text(slide, x, y, w, 0.28, "Accuracy global final promedio", 9.7, C.INK, True, PP_ALIGN.CENTER)
    start = 0.88
    span = 0.12
    row_h = (h - 0.35) / len(data)
    for i, (label, value, color) in enumerate(data):
        yy = y + 0.42 + i * row_h
        add_text(slide, x, yy + 0.05, 2.45, 0.25, label, 8.7, C.INK, True)
        add_rect(slide, x + 2.65, yy + 0.07, w - 4.05, 0.20, RGBColor(230, 233, 236), RGBColor(230, 233, 236))
        bw = (w - 4.05) * max(0, min(1, (value - start) / span))
        add_rect(slide, x + 2.65, yy + 0.07, bw, 0.20, color, color)
        add_text(slide, x + w - 1.20, yy - 0.01, 1.15, 0.30, f"{value*100:.2f}%", 9.0, color, True, PP_ALIGN.RIGHT)


def build_pptx():
    extract_template_assets()
    prs = Presentation()
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C.PAPER, C.PAPER)
    slide.shapes.add_picture(str(HEADER_ORANGE), inch(0), inch(0), width=inch(SLIDE_W), height=inch(2.72))

    # Committee/advisor strip.
    add_rect(slide, 0.42, 2.72, 13.55, 0.55, C.PEACH, C.DASH, dashed=True)
    add_rect(slide, 13.97, 2.72, 13.17, 0.55, C.PEACH, C.DASH, dashed=True)
    add_text(slide, 0.55, 2.84, 13.1, 0.24, "COMITÉ - Seguridad, Redes e IoT", 9.4, C.INK, True, PP_ALIGN.CENTER)
    add_text(slide, 14.10, 2.84, 12.85, 0.24, "Asesor: Prof. Carlos Andrés Lozano Garzón, PhD", 8.8, C.INK, False, PP_ALIGN.RIGHT)

    # Title and authors block.
    add_rect(slide, 0.75, 3.35, 26.05, 3.1, C.WHITE, C.DASH, dashed=True)
    add_text(
        slide,
        1.2,
        3.62,
        25.1,
        0.98,
        "Sistema de detección de intrusiones para redes IoT/MQTT mediante TinyML, HFL y ASCON-128",
        22.0,
        C.INK,
        False,
        PP_ALIGN.CENTER,
    )
    add_rect(slide, 0.75, 4.72, 26.05, 0.01, C.WHITE, C.DASH, dashed=True)
    add_text(slide, 1.2, 4.98, 25.1, 0.36, "Santiago Alejandro Jaimes Puerto  |  Nicolás Casas Ibarra", 11.7, C.BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, 1.2, 5.46, 25.1, 0.34, "Correos institucionales: por completar", 8.7, C.BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, 1.2, 5.82, 25.1, 0.32, "Universidad de los Andes - Departamento de Ingeniería de Sistemas y Computación - 2026", 8.2, C.MUTED, False, PP_ALIGN.CENTER)

    # Summary full-width section.
    sx, sy, sw, sh = add_section(slide, 0.75, 6.75, 26.05, 3.4, "Resumen", "blue")
    add_text(
        slide,
        sx,
        sy,
        sw,
        sh,
        "Esta tesis diseña, implementa y evalúa un IDS para tráfico IoT/MQTT sobre una arquitectura Edge-Fog-Cloud. "
        "El sistema ejecuta inferencia TinyML en ESP32-S3, entrenamiento local y agregación en Raspberry Pi, y coordinación global en PC. "
        "Se comparan variantes RN/MLP, CNN-1D, ASCON/no-ASCON y FOG/no-FOG. Los resultados muestran que HFL puede operar en hardware real con accuracy global entre 92.98% y 97.50%, "
        "mientras ASCON-128 protege pesos y features con overhead temporal menor a 0.1% del ciclo federado.",
        12.0,
    )

    left_x, right_x = 0.75, 14.05
    col_w = 12.75
    y0 = 10.65

    # Left column: Introduction, objectives, methodology.
    x, y, w, h = add_section(slide, left_x, y0, col_w, 4.15, "Introducción", "orange")
    add_bullets(
        slide,
        x,
        y,
        w,
        h,
        [
            "Las redes IoT/IIoT aumentan la superficie de ataque y usan hardware restringido.",
            "Los IDS centralizados exigen mover datos crudos y modelos grandes.",
            "FL reduce transferencia de datos, pero los pesos y métricas siguen siendo activos sensibles.",
            "ASCON-128 se integra como cifrado autenticado ligero para proteger el ciclo federado.",
        ],
        11.3,
    )

    x, y, w, h = add_section(slide, left_x, 15.2, col_w, 4.45, "Objetivos", "orange")
    add_text(slide, x, y, w, 0.42, "Objetivo general", 11.6, C.INK, True)
    add_text(slide, x, y + 0.48, w, 0.72, "Diseñar, implementar y evaluar un IDS IoT/MQTT basado en TinyML, HFL y ASCON-128.", 10.2)
    add_bullets(
        slide,
        x,
        y + 1.38,
        w,
        h - 1.4,
        [
            "Construir pipeline de 13 features y tres clases.",
            "Comparar modelos compactos: MLP/RN, CNN-1D, Residual-MLP y Transformer.",
            "Implementar HFL bidireccional PC-RPi-ESP32.",
            "Medir impacto de ASCON, no-ASCON y FOG.",
        ],
        10.0,
    )

    x, y, w, h = add_section(slide, left_x, 20.05, col_w, 11.7, "Metodología", "orange")
    add_bullets(
        slide,
        x,
        y,
        w,
        2.25,
        [
            "Fase 1: preparación de dataset con 256 276 flujos y StandardScaler.",
            "Fase 2: entrenamiento offline y exportación de pesos/scaler.",
            "Fase 3: despliegue TinyML en ESP32-S3.",
            "Fase 4: HFL con Raspberry Pi y PC usando FedAvg ponderado.",
            "Fase 5: cifrado ASCON-128 y baseline JSON plano.",
            "Fase 6: análisis temporal de accuracy, loss y magnitud de pesos.",
        ],
        10.0,
    )
    add_rect(slide, x, y + 2.65, w, 2.05, RGBColor(245, 247, 248), C.DASH, dashed=True)
    add_text(slide, x + 0.18, y + 2.83, w - 0.36, 0.32, "Configuración HFL", 10.2, C.BLUE, True)
    add_text(
        slide,
        x + 0.18,
        y + 3.25,
        w - 0.36,
        1.1,
        "Gateway: 40 muestras por actualización, 5 épocas, batch=8, Adam lr=0.005.\n"
        "MLP/RN: 13 -> 32 -> 16 -> 8 -> 3; 1 139 parámetros; 163 federados.\n"
        "CNN-1D: Conv1D + BN + GAP + densas federadas.",
        9.0,
    )
    add_rect(slide, x, y + 5.05, w, 4.85, RGBColor(245, 247, 248), C.DASH, dashed=True)
    add_text(slide, x + 0.18, y + 5.25, w - 0.36, 0.32, "Variantes evaluadas", 10.2, C.BLUE, True)
    add_text(
        slide,
        x + 0.18,
        y + 5.67,
        w - 0.36,
        3.75,
        "• hfl_v7-RN: MLP/RN + ASCON, estándar y FOG.\n"
        "• hfl_v7-no-ascon: MLP/RN + JSON plano, estándar y FOG.\n"
        "• hfl_v7-CNN: CNN-1D + ASCON, estándar y FOG.\n"
        "• FOG: leader/peer con topics fog/weights y fog/global_model.\n"
        "• No-ASCON: topics fl/features_plain y fl/global_model_plain.",
        9.4,
    )

    # Right column: Design, results, conclusions.
    x, y, w, h = add_section(slide, right_x, y0, col_w, 8.75, "Diseño e implementación", "orange")
    mini_architecture(slide, x, y, w, 4.8)
    add_bullets(
        slide,
        x,
        y + 4.85,
        w,
        h - 4.9,
        [
            "PC: FastAPI, modelo global, dashboard y FedAvg global.",
            "Raspberry Pi: broker MQTT, etiquetado heurístico, entrenamiento local y FedAvg fog.",
            "ESP32-S3: inferencia TinyML, cálculo/publicación de features y alertamiento.",
            "ASCON: payload autenticado {ct, tag, nonce}; no-ASCON: JSON plano comparable.",
        ],
        9.5,
    )

    x, y, w, h = add_section(slide, right_x, 19.8, col_w, 12.25, "Resultados", "orange")
    small_results_table(slide, x, y, w, 2.05)
    accuracy_bars(slide, x, y + 2.35, w, 1.95)
    add_text(slide, x, y + 4.48, w, 0.34, "Curvas federadas v7: accuracy/loss y magnitud de pesos", 9.2, C.INK, True, PP_ALIGN.CENTER)
    gap = 0.22
    img_w = (w - gap) / 2
    add_image_fit(slide, ACC_LOSS_IMG, x, y + 4.95, img_w, 4.15)
    add_image_fit(slide, WEIGHTS_IMG, x + img_w + gap, y + 4.95, img_w, 4.15)
    add_text(
        slide,
        x,
        y + 9.35,
        w,
        0.95,
        "ASCON añade ~36% de tamaño por mensaje, pero el costo temporal medido es ~42 ms por ronda, equivalente a ~0.06% del ciclo FL.",
        9.2,
        C.MUTED,
    )

    x, y, w, h = add_section(slide, right_x, 32.55, col_w, 4.45, "Conclusiones", "blue")
    add_bullets(
        slide,
        x,
        y,
        w,
        h,
        [
            "La arquitectura HFL bidireccional fue validada en hardware real ESP32-S3, Raspberry Pi y PC.",
            "El MLP/RN ofrece el mejor balance entre simplicidad edge, memoria e inferencia.",
            "CNN-1D+FOG obtuvo el mejor accuracy promedio, aunque con mayor complejidad de implementación.",
            "ASCON-128 protege el ciclo federado con overhead temporal despreciable.",
            "Trabajo futuro: más gateways, datos non-IID reales y pruebas adversariales de poisoning.",
        ],
        9.4,
    )

    # Footer.
    add_rect(slide, 0, 37.85, SLIDE_W, 0.9, C.PEACH, C.PEACH)
    add_text(slide, 3.7, 38.13, 20.2, 0.32, "http://sistemas.uniandes.edu.co     /DISCuniandes     @ISISUniandes", 8.5, C.INK, True, PP_ALIGN.CENTER)

    prs.save(OUT_PPTX)


# ---------------------------------------------------------------------------
# PNG preview, built from the same template-like layout.


def font(size: int, bold=False):
    candidates = [
        r"C:\Windows\Fonts\segoeuib.ttf" if bold else r"C:\Windows\Fonts\segoeui.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def wrap_draw(draw, xy, text, width, size, fill, bold=False, spacing=1.15):
    x, y = xy
    fnt = font(size, bold)
    for raw in text.split("\n"):
        words = raw.split()
        line = ""
        if not words:
            y += int(size * spacing)
            continue
        for word in words:
            probe = f"{line} {word}".strip()
            if draw.textbbox((0, 0), probe, font=fnt)[2] <= width:
                line = probe
            else:
                draw.text((x, y), line, font=fnt, fill=fill)
                y += int(size * spacing)
                line = word
        if line:
            draw.text((x, y), line, font=fnt, fill=fill)
            y += int(size * spacing)
    return y


def build_preview():
    extract_template_assets()
    scale = 82
    W, H = int(SLIDE_W * scale), int(SLIDE_H * scale)
    img = Image.new("RGB", (W, H), rgb_hex(C.PAPER))
    draw = ImageDraw.Draw(img)

    def sx(v): return int(v * scale)

    header = Image.open(HEADER_ORANGE).convert("RGB").resize((W, sx(2.72)), Image.Resampling.LANCZOS)
    img.paste(header, (0, 0))

    def dashed_rect(x, y, w, h, fill="white"):
        draw.rectangle([sx(x), sx(y), sx(x + w), sx(y + h)], fill=fill, outline=rgb_hex(C.DASH), width=2)

    def bar(x, y, w, title, color="orange"):
        fill = rgb_hex(C.BLUE if color == "blue" else C.ORANGE)
        draw.rectangle([sx(x), sx(y), sx(x + w), sx(y + 0.66)], fill=fill)
        draw.text((sx(x + 0.28), sx(y + 0.14)), title, font=font(22, True), fill="white")

    # Top strip and title.
    dashed_rect(0.42, 2.72, 13.55, 0.55, rgb_hex(C.PEACH))
    dashed_rect(13.97, 2.72, 13.17, 0.55, rgb_hex(C.PEACH))
    draw.text((sx(3.05), sx(2.86)), "COMITÉ - Seguridad, Redes e IoT", font=font(14, True), fill=rgb_hex(C.INK))
    draw.text((sx(17.1), sx(2.86)), "Asesor: Prof. Carlos Andrés Lozano Garzón, PhD", font=font(13), fill=rgb_hex(C.INK))
    dashed_rect(0.75, 3.35, 26.05, 3.1)
    wrap_draw(draw, (sx(1.9), sx(3.62)), "Sistema de detección de intrusiones para redes IoT/MQTT mediante TinyML, HFL y ASCON-128", sx(23.8), 34, rgb_hex(C.INK))
    draw.text((sx(8.15), sx(5.02)), "Santiago Alejandro Jaimes Puerto  |  Nicolás Casas Ibarra", font=font(18, True), fill=rgb_hex(C.BLUE))
    draw.text((sx(10.2), sx(5.52)), "Correos institucionales: por completar", font=font(14, True), fill=rgb_hex(C.BLUE))
    draw.text((sx(6.95), sx(5.94)), "Universidad de los Andes - Departamento de Ingeniería de Sistemas y Computación - 2026", font=font(13), fill=rgb_hex(C.MUTED))

    # Summary.
    bar(0.75, 6.75, 26.05, "Resumen", "blue")
    dashed_rect(0.75, 7.41, 26.05, 2.74)
    wrap_draw(draw, (sx(1.12), sx(7.75)), "Esta tesis diseña, implementa y evalúa un IDS para tráfico IoT/MQTT sobre Edge-Fog-Cloud. Ejecuta TinyML en ESP32-S3, entrenamiento local en Raspberry Pi y coordinación global en PC. Compara RN/MLP, CNN-1D, ASCON/no-ASCON y FOG/no-FOG; HFL logra 92.98%-97.50% de accuracy global y ASCON protege el ciclo federado con overhead temporal menor a 0.1%.", sx(25.3), 18, rgb_hex(C.INK))

    # Sections.
    sections = [
        (0.75, 10.65, 12.75, 4.15, "Introducción", "orange", "• IoT/IIoT aumenta superficie de ataque.\n• Los IDS centralizados mueven datos crudos y requieren modelos grandes.\n• FL reduce transferencia, pero pesos y métricas siguen siendo sensibles.\n• ASCON-128 protege el ciclo federado."),
        (0.75, 15.2, 12.75, 4.45, "Objetivos", "orange", "General: diseñar, implementar y evaluar un IDS IoT/MQTT con TinyML, HFL y ASCON-128.\n\n• Pipeline de 13 features y 3 clases.\n• Comparar modelos compactos.\n• Implementar HFL PC-RPi-ESP32.\n• Medir ASCON/no-ASCON/FOG."),
        (0.75, 20.05, 12.75, 11.7, "Metodología", "orange", "• Dataset: 256 276 flujos, StandardScaler.\n• Entrenamiento offline y exportación de pesos.\n• Despliegue TinyML en ESP32-S3.\n• Gateway: 40 muestras, 5 épocas, batch=8.\n• PC: FedAvg global y dashboard.\n• Variantes: hfl_v7-RN, hfl_v7-no-ascon, hfl_v7-CNN."),
        (14.05, 10.65, 12.75, 8.75, "Diseño e implementación", "orange", "PC Cloud + FedAvg\n    modelo global / pesos\nRaspberry Pi FOG leader <-> Raspberry Pi FOG peer\n    MQTT\nESP32-S3 Edge A        ESP32-S3 Edge B\n\nASCON: {ct, tag, nonce}; no-ASCON: JSON plano *_plain.\nFOG: leader preagrega pesos antes del PC."),
        (14.05, 19.8, 12.75, 11.85, "Resultados", "orange", "RN+ASCON: 94.63% | loss 0.1415 | 61.25 s\nRN sin ASCON: 92.98% | loss 0.1581 | 51.01 s\nCNN-1D+FOG: 97.50% | loss 0.0910 | 58.61 s"),
        (14.05, 32.05, 12.75, 4.95, "Conclusiones", "blue", "• HFL bidireccional validado en hardware real.\n• MLP/RN balancea memoria, simplicidad y desempeño.\n• CNN-1D+FOG logra mayor accuracy promedio.\n• ASCON aporta seguridad con overhead temporal despreciable.\n• Futuro: más gateways, non-IID real y pruebas de poisoning."),
    ]
    for x, y, w, h, title, color, text in sections:
        bar(x, y, w, title, color)
        dashed_rect(x, y + 0.66, w, h - 0.66)
        body_size = 17 if title in {"Introducción", "Objetivos", "Conclusiones"} else 16
        if title in {"Metodología", "Diseño e implementación"}:
            body_size = 15
        wrap_draw(draw, (sx(x + 0.38), sx(y + 1.02)), text, sx(w - 0.76), body_size, rgb_hex(C.INK))

    # Result images in preview.
    def paste_fit(path, x, y, w, h):
        im = Image.open(path).convert("RGB")
        im.thumbnail((sx(w), sx(h)), Image.Resampling.LANCZOS)
        img.paste(im, (sx(x) + (sx(w) - im.width) // 2, sx(y) + (sx(h) - im.height) // 2))

    # Extra visible result emphasis in the preview.
    bx, by, bw = 14.42, 22.0, 12.0
    rows = [("RN + ASCON", 0.9463, C.BLUE), ("RN sin ASCON", 0.9298, RGBColor(45, 155, 128)), ("CNN-1D + FOG", 0.9750, C.ORANGE)]
    draw.text((sx(bx), sx(by)), "Accuracy global final promedio", font=font(15, True), fill=rgb_hex(C.INK))
    for i, (label, val, color) in enumerate(rows):
        yy = by + 0.42 + i * 0.38
        draw.text((sx(bx), sx(yy)), label, font=font(13, True), fill=rgb_hex(C.INK))
        draw.rectangle([sx(bx + 2.5), sx(yy + 0.05), sx(bx + 10.5), sx(yy + 0.23)], fill="#e6e9ec")
        width = 8.0 * max(0, min(1, (val - 0.88) / 0.12))
        draw.rectangle([sx(bx + 2.5), sx(yy + 0.05), sx(bx + 2.5 + width), sx(yy + 0.23)], fill=rgb_hex(color))
        draw.text((sx(bx + 10.75), sx(yy - 0.02)), f"{val*100:.2f}%", font=font(13, True), fill=rgb_hex(color))
    paste_fit(ACC_LOSS_IMG, 14.42, 24.75, 5.85, 4.0)
    paste_fit(WEIGHTS_IMG, 20.57, 24.75, 5.85, 4.0)

    # Footer.
    draw.rectangle([0, sx(37.85), W, sx(38.75)], fill=rgb_hex(C.PEACH))
    draw.text((sx(8.3), sx(38.16)), "http://sistemas.uniandes.edu.co     /DISCuniandes     @ISISUniandes", font=font(12, True), fill=rgb_hex(C.INK))

    img.save(OUT_PREVIEW)


if __name__ == "__main__":
    build_pptx()
    build_preview()
    print(f"PPTX: {OUT_PPTX}")
    print(f"Preview: {OUT_PREVIEW}")
