# -*- coding: utf-8 -*-
"""Build the thesis poster for the HFL v7 project.

The poster is generated as an editable PPTX and a PNG preview. It uses the
Uniandes poster template header supplied by the user, but rebuilds the content
as native PowerPoint shapes/text so it remains easy to edit.
"""

from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile
import textwrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SRC = ROOT.parents[1]
TEMPLATE = Path(r"C:\Users\VivoBook\Downloads\FormatoPoster_ISIS_v2.potx")
ASSET_DIR = ROOT / "images" / "Poster"
OUT_PPTX = ROOT / "Poster_Tesis_HFL_v7.pptx"
OUT_PREVIEW = ROOT / "Poster_Tesis_HFL_v7_preview.png"

HEADER_IMG = ASSET_DIR / "uniandes_header_green.png"
ACC_LOSS_IMG = ROOT / "images" / "Resultados" / "hfl_v7_accuracy_loss_matrix.png"
WEIGHTS_IMG = ROOT / "images" / "Resultados" / "hfl_v7_weight_magnitude_matrix.png"

EMU_PER_IN = 914400

SLIDE_W = 27.56
SLIDE_H = 39.38

COLORS = {
    "green": RGBColor(0, 105, 55),
    "green_dark": RGBColor(0, 74, 44),
    "green_mid": RGBColor(35, 137, 88),
    "green_light": RGBColor(230, 244, 236),
    "blue": RGBColor(18, 96, 163),
    "blue_light": RGBColor(232, 242, 252),
    "orange": RGBColor(230, 126, 34),
    "orange_light": RGBColor(254, 241, 229),
    "red": RGBColor(190, 54, 54),
    "ink": RGBColor(25, 35, 45),
    "muted": RGBColor(83, 94, 104),
    "line": RGBColor(204, 214, 220),
    "paper": RGBColor(248, 250, 248),
    "white": RGBColor(255, 255, 255),
}


def inch(value: float) -> int:
    return int(value * EMU_PER_IN)


def rgb_hex(rgb: RGBColor) -> str:
    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"


def extract_header() -> None:
    """Extract the selected green institutional header from the POTX."""
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    if HEADER_IMG.exists():
        return
    with ZipFile(TEMPLATE) as zf:
        # This is the green "Facultad de Ingeniería" header in the template.
        data = zf.read("ppt/media/image32.png")
    HEADER_IMG.write_bytes(data)


def set_textbox_text(
    shape,
    text: str,
    size: float,
    color: RGBColor = COLORS["ink"],
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Segoe UI",
    valign=MSO_ANCHOR.TOP,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.03)
    tf.margin_right = inch(0.03)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    tf.vertical_anchor = valign
    for idx, raw_line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = raw_line
        p.alignment = align
        p.space_after = Pt(0)
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float,
    color: RGBColor = COLORS["ink"],
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Segoe UI",
    valign=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    set_textbox_text(tb, text, size, color, bold, align, font, valign)
    return tb


def add_shape(slide, shape_type, x: float, y: float, w: float, h: float, fill, line=None):
    shp = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.color.rgb = fill
    else:
        shp.line.color.rgb = line
    return shp


def add_card(slide, x: float, y: float, w: float, h: float, title: str, accent=COLORS["green"]):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, COLORS["white"], COLORS["line"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.28, accent, accent)
    add_text(slide, x + 0.25, y + 0.36, w - 0.5, 0.42, title.upper(), 15, accent, True)
    return x + 0.25, y + 0.88, w - 0.5, h - 1.05


def add_metric(slide, x: float, y: float, w: float, label: str, value: str, note: str, accent):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 1.65, COLORS["white"], COLORS["line"])
    add_text(slide, x + 0.2, y + 0.18, w - 0.4, 0.3, label.upper(), 9.5, COLORS["muted"], True, PP_ALIGN.CENTER)
    add_text(slide, x + 0.15, y + 0.47, w - 0.3, 0.58, value, 23, accent, True, PP_ALIGN.CENTER)
    add_text(slide, x + 0.18, y + 1.11, w - 0.36, 0.35, note, 8.8, COLORS["muted"], False, PP_ALIGN.CENTER)


def bullet_text(slide, x: float, y: float, w: float, h: float, lines: list[str], size=12.3):
    text = "\n".join(f"• {line}" for line in lines)
    return add_text(slide, x, y, w, h, text, size, COLORS["ink"])


def add_result_table(slide, x: float, y: float, w: float, h: float):
    rows = [
        ("Variante", "Acc.", "Loss", "Ronda"),
        ("RN + ASCON", "94.63%", "0.1415", "61.25 s"),
        ("RN sin ASCON", "92.98%", "0.1581", "51.01 s"),
        ("CNN-1D + FOG", "97.50%", "0.0910", "58.61 s"),
    ]
    col_w = [0.43, 0.19, 0.18, 0.20]
    row_h = h / len(rows)
    for r, row in enumerate(rows):
        yy = y + r * row_h
        fill = COLORS["green"] if r == 0 else (COLORS["green_light"] if r % 2 == 1 else COLORS["white"])
        text_color = COLORS["white"] if r == 0 else COLORS["ink"]
        xx = x
        for c, val in enumerate(row):
            cw = w * col_w[c]
            add_shape(slide, MSO_SHAPE.RECTANGLE, xx, yy, cw, row_h, fill, COLORS["line"])
            add_text(slide, xx + 0.08, yy + 0.08, cw - 0.16, row_h - 0.12, val, 9.3 if r else 9.8, text_color, r == 0)
            xx += cw


def add_accuracy_chart(slide, x: float, y: float, w: float, h: float):
    chart_data = CategoryChartData()
    chart_data.categories = ["RN+ASCON", "RN sin\nASCON", "CNN FOG"]
    chart_data.add_series("Accuracy global final", (0.9463, 0.9298, 0.9750))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, inch(x), inch(y), inch(w), inch(h), chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0.88
    chart.value_axis.maximum_scale = 1.0
    chart.value_axis.major_unit = 0.02
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = COLORS["green_mid"]
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Accuracy global final promedio"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True


def add_architecture_diagram(slide, x: float, y: float, w: float, h: float):
    add_text(slide, x, y, w, 0.35, "Flujo bidireccional Edge-Fog-Cloud", 12.5, COLORS["green_dark"], True)

    pc = (x + w * 0.35, y + 0.75, w * 0.30, 0.78)
    leader = (x + w * 0.13, y + 2.25, w * 0.33, 0.78)
    peer = (x + w * 0.54, y + 2.25, w * 0.33, 0.78)
    edge_a = (x + w * 0.12, y + 3.85, w * 0.34, 0.78)
    edge_b = (x + w * 0.54, y + 3.85, w * 0.34, 0.78)

    for box, label, fill in [
        (pc, "PC Cloud\nFedAvg global", COLORS["blue_light"]),
        (leader, "RPi FOG Leader\nMQTT + train + FedAvg", COLORS["green_light"]),
        (peer, "RPi FOG Peer\ntrain local", COLORS["green_light"]),
        (edge_a, "ESP32-S3 Edge A\nTinyML + features", COLORS["orange_light"]),
        (edge_b, "ESP32-S3 Edge B\nTinyML + features", COLORS["orange_light"]),
    ]:
        bx, by, bw, bh = box
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh, fill, COLORS["line"])
        add_text(slide, bx + 0.08, by + 0.1, bw - 0.16, bh - 0.1, label, 8.8, COLORS["ink"], True, PP_ALIGN.CENTER)

    def down_arrow(ax, ay, aw, ah, color, label):
        add_shape(slide, MSO_SHAPE.DOWN_ARROW, ax, ay, aw, ah, color, color)
        add_text(slide, ax - 0.7, ay + ah * 0.25, aw + 1.4, 0.24, label, 6.8, COLORS["muted"], False, PP_ALIGN.CENTER)

    def up_arrow(ax, ay, aw, ah, color, label):
        add_shape(slide, MSO_SHAPE.UP_ARROW, ax, ay, aw, ah, color, color)
        add_text(slide, ax - 0.8, ay + ah * 0.38, aw + 1.6, 0.24, label, 6.8, COLORS["muted"], False, PP_ALIGN.CENTER)

    down_arrow(x + w * 0.49, y + 1.55, 0.32, 0.55, COLORS["green_mid"], "modelo global")
    down_arrow(x + w * 0.25, y + 3.08, 0.32, 0.55, COLORS["green_mid"], "fl/global_model")
    down_arrow(x + w * 0.67, y + 3.08, 0.32, 0.55, COLORS["green_mid"], "fl/global_model")
    up_arrow(x + w * 0.34, y + 3.08, 0.28, 0.55, COLORS["orange"], "features")
    up_arrow(x + w * 0.76, y + 3.08, 0.28, 0.55, COLORS["orange"], "features")
    up_arrow(x + w * 0.43, y + 1.55, 0.28, 0.55, COLORS["orange"], "pesos")

    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x + w * 0.47, y + 2.43, w * 0.06, 0.35, COLORS["green_mid"], COLORS["green_mid"])
    add_text(slide, x + w * 0.45, y + 2.03, w * 0.12, 0.24, "fog/global_model", 6.7, COLORS["muted"], False, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.LEFT_ARROW, x + w * 0.47, y + 2.78, w * 0.06, 0.35, COLORS["orange"], COLORS["orange"])
    add_text(slide, x + w * 0.45, y + 3.18, w * 0.12, 0.24, "fog/weights", 6.7, COLORS["muted"], False, PP_ALIGN.CENTER)

    add_text(
        slide,
        x,
        y + h - 0.72,
        w,
        0.55,
        "ASCON: payload {ct, tag, nonce}. Sin ASCON: JSON plano y topics *_plain.",
        8.7,
        COLORS["muted"],
        False,
        PP_ALIGN.CENTER,
    )


def add_image_fit(slide, image_path: Path, x: float, y: float, w: float, h: float):
    img = Image.open(image_path)
    iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        width = w
        height = w / img_ratio
        yy = y + (h - height) / 2
        xx = x
    else:
        height = h
        width = h * img_ratio
        xx = x + (w - width) / 2
        yy = y
    return slide.shapes.add_picture(str(image_path), inch(xx), inch(yy), width=inch(width), height=inch(height))


def build_pptx() -> None:
    extract_header()

    prs = Presentation()
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background and institutional header.
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, COLORS["paper"], COLORS["paper"])
    slide.shapes.add_picture(str(HEADER_IMG), inch(0), inch(0), width=inch(SLIDE_W), height=inch(2.72))

    add_text(
        slide,
        0.85,
        3.05,
        19.8,
        0.78,
        "Sistema de detección de intrusiones para redes IoT/MQTT",
        29,
        COLORS["green_dark"],
        True,
    )
    add_text(
        slide,
        0.85,
        3.9,
        18.7,
        0.48,
        "TinyML en ESP32-S3 + Aprendizaje Federado Jerárquico + ASCON-128 en arquitectura Edge-Fog-Cloud",
        15,
        COLORS["ink"],
        False,
    )
    add_text(
        slide,
        20.4,
        3.22,
        6.2,
        0.95,
        "Santiago Alejandro Jaimes Puerto\nNicolás Casas Ibarra\nDirector: Prof. Carlos Andrés Lozano Garzón, PhD",
        10.8,
        COLORS["muted"],
        False,
        PP_ALIGN.RIGHT,
    )

    metrics = [
        ("Features", "13", "flujos MQTT/TCP"),
        ("Clases", "3", "normal, brute force, scan"),
        ("Modelo Edge", "4.5 KB", "MLP/RN en ESP32"),
        ("FedAvg", "163", "parámetros federados"),
        ("Ronda local", "40", "muestras + 5 épocas"),
        ("ASCON", "0.06%", "overhead temporal"),
    ]
    mx, my, mw, mgap = 0.85, 5.05, 4.05, 0.28
    for i, (label, value, note) in enumerate(metrics):
        add_metric(slide, mx + i * (mw + mgap), my, mw, label, value, note, COLORS["green_mid"] if i != 5 else COLORS["orange"])

    # Columns.
    col_y = 7.25
    left_x, mid_x, right_x = 0.85, 9.55, 18.25
    col_w, right_w = 8.25, 8.45

    # Left column.
    ix, iy, iw, ih = add_card(slide, left_x, col_y, col_w, 4.35, "Problema y objetivo", COLORS["green"])
    add_text(
        slide,
        ix,
        iy,
        iw,
        ih,
        "Los IDS centralizados alcanzan alta precisión, pero exigen memoria, conectividad y centralización de datos que no son realistas en IoT.\n\nObjetivo: validar un IDS distribuido que detecte ataques MQTT en el borde, entrene sin mover datos crudos y proteja pesos/features durante el ciclo federado.",
        11.3,
    )

    ix, iy, iw, ih = add_card(slide, left_x, 12.05, col_w, 5.35, "Pipeline experimental", COLORS["blue"])
    bullet_text(
        slide,
        ix,
        iy,
        iw,
        ih,
        [
            "Dataset de 256 276 flujos con 13 features normalizadas.",
            "Problema triclase: normal, mqtt_bruteforce y scan_A.",
            "ESP32-S3 ejecuta inferencia TinyML y publica features por MQTT.",
            "Raspberry Pi etiqueta heurísticamente, acumula muestras y entrena localmente.",
            "PC coordina rondas, FedAvg global y dashboard de métricas.",
        ],
        10.5,
    )

    ix, iy, iw, ih = add_card(slide, left_x, 17.9, col_w, 5.75, "Modelos evaluados", COLORS["green"])
    add_text(slide, ix, iy, iw, 0.52, "MLP/RN principal", 12.2, COLORS["green_dark"], True)
    bullet_text(
        slide,
        ix,
        iy + 0.55,
        iw,
        1.65,
        [
            "Arquitectura 13 -> 32 -> 16 -> 8 -> 3.",
            "1 139 parámetros; aprox. 4.5 KB en float32.",
            "Capas federadas: W3, b3, W4, b4.",
        ],
        9.8,
    )
    add_text(slide, ix, iy + 2.42, iw, 0.52, "CNN-1D experimental", 12.2, COLORS["green_dark"], True)
    bullet_text(
        slide,
        ix,
        iy + 2.98,
        iw,
        1.65,
        [
            "Entrada (13,1), Conv1D + BatchNorm + GAP.",
            "Solo se federan las capas densas finales.",
            "Más expresiva, pero más costosa que MLP en C++.",
        ],
        9.8,
    )

    ix, iy, iw, ih = add_card(slide, left_x, 24.15, col_w, 5.45, "Seguridad y variantes", COLORS["orange"])
    bullet_text(
        slide,
        ix,
        iy,
        iw,
        ih,
        [
            "ASCON-128 protege confidencialidad e integridad con payload {ct, tag, nonce}.",
            "La rama no-ASCON usa JSON plano y topics/endpoints *_plain.",
            "El modo FOG agrega leader/peer antes del PC para reducir tráfico hacia Cloud.",
            "Se comparan RN+ASCON, RN sin ASCON, CNN-1D y CNN-1D+FOG.",
        ],
        10.4,
    )

    ix, iy, iw, ih = add_card(slide, left_x, 30.1, col_w, 6.55, "Aporte de la tesis", COLORS["green"])
    bullet_text(
        slide,
        ix,
        iy,
        iw,
        ih,
        [
            "Demuestra un ciclo HFL bidireccional completo en hardware real.",
            "Integra TinyML, MQTT, FastAPI, FedAvg y criptografía ligera.",
            "Mide convergencia, pérdida, magnitud de pesos y overhead de seguridad.",
            "Aísla tres efectos: cifrado, modelo y topología FOG.",
        ],
        10.8,
    )

    # Middle column.
    ix, iy, iw, ih = add_card(slide, mid_x, col_y, col_w, 7.85, "Arquitectura implementada", COLORS["green"])
    add_architecture_diagram(slide, ix, iy, iw, ih)

    ix, iy, iw, ih = add_card(slide, mid_x, 15.6, col_w, 6.4, "Flujo bidireccional v7", COLORS["blue"])
    add_text(
        slide,
        ix,
        iy,
        iw,
        ih,
        "1. PC despliega modelo global a Raspberry Pi.\n"
        "2. Gateway publica modelo a ESP32 por MQTT.\n"
        "3. ESP32 infiere localmente y publica features.\n"
        "4. Gateway etiqueta, entrena con 40 muestras y calcula pesos locales.\n"
        "5. Sin FOG: gateway reporta al PC.\n"
        "6. Con FOG: peer reporta a leader, leader preagrega y reporta al PC.\n"
        "7. PC actualiza modelo global y reinicia el ciclo.",
        10.4,
    )

    ix, iy, iw, ih = add_card(slide, mid_x, 22.55, col_w, 5.8, "Contraste de rutas", COLORS["orange"])
    add_result_table(slide, ix, iy, iw, 2.55)
    add_text(
        slide,
        ix,
        iy + 2.9,
        iw,
        1.65,
        "Con ASCON se mantienen topics base (`fl/features`, `fl/global_model`) y el contenido viaja cifrado. Sin ASCON se usan rutas *_plain para medir el baseline sin costo criptográfico.",
        10.1,
        COLORS["ink"],
    )

    ix, iy, iw, ih = add_card(slide, mid_x, 28.85, col_w, 7.8, "Hallazgos de operación", COLORS["green"])
    bullet_text(
        slide,
        ix,
        iy,
        iw,
        ih,
        [
            "La convergencia HFL se estabiliza alrededor de 90-95% de accuracy.",
            "El cifrado no degrada de forma observable la convergencia.",
            "FOG permite preagregar gateways antes del PC y escalar la topología.",
            "CNN-1D+FOG obtiene el mayor accuracy observado, pero MLP/RN sigue siendo la opción más simple para ESP32.",
            "Comunicación por ronda: aprox. 7.2 KB en JSON plano y 9.8 KB con ASCON.",
        ],
        10.7,
    )

    # Right column.
    ix, iy, iw, ih = add_card(slide, right_x, col_y, right_w, 5.2, "Resultados HFL v7", COLORS["green"])
    add_accuracy_chart(slide, ix, iy + 0.05, iw, 2.45)
    add_text(
        slide,
        ix,
        iy + 2.72,
        iw,
        1.05,
        "RN+ASCON mantiene desempeño competitivo frente al baseline sin cifrado. CNN-1D+FOG presenta el mejor resultado promedio en las corridas analizadas.",
        9.8,
        COLORS["muted"],
    )

    ix, iy, iw, ih = add_card(slide, right_x, 12.95, right_w, 8.05, "Evolución de accuracy y loss", COLORS["blue"])
    add_image_fit(slide, ACC_LOSS_IMG, ix, iy, iw, ih - 0.15)

    ix, iy, iw, ih = add_card(slide, right_x, 21.55, right_w, 8.05, "Evolución de magnitud de pesos", COLORS["orange"])
    add_image_fit(slide, WEIGHTS_IMG, ix, iy, iw, ih - 0.15)

    ix, iy, iw, ih = add_card(slide, right_x, 30.15, right_w, 6.5, "Conclusión", COLORS["green"])
    add_text(
        slide,
        ix,
        iy,
        iw,
        ih,
        "La combinación de TinyML, HFL y ASCON-128 es viable para IDS IoT/MQTT en una arquitectura Edge-Fog-Cloud. El sistema evita centralizar datos crudos, mantiene actualizaciones bidireccionales de modelo, agrega seguridad autenticada con costo temporal despreciable y permite comparar el efecto de modelo, cifrado y topología.",
        12.4,
        COLORS["ink"],
        True,
    )

    # Footer.
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 37.85, SLIDE_W, 1.0, COLORS["green_dark"], COLORS["green_dark"])
    add_text(
        slide,
        0.85,
        38.08,
        18.0,
        0.38,
        "Artefactos: src/hfl_v7-RN | src/hfl_v7-no-ascon | src/hfl_v7-CNN | Documento_Final | Analisis de Modelos",
        9.5,
        COLORS["white"],
    )
    add_text(
        slide,
        19.0,
        38.08,
        7.7,
        0.38,
        "Universidad de los Andes - Ingeniería de Sistemas y Computación - 2026",
        9.5,
        COLORS["white"],
        False,
        PP_ALIGN.RIGHT,
    )

    prs.save(OUT_PPTX)


# ---------------------------------------------------------------------------
# PNG preview


def font(size: int, bold: bool = False):
    candidates = [
        r"C:\Windows\Fonts\segoeuib.ttf" if bold else r"C:\Windows\Fonts\segoeui.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def draw_wrapped(draw, xy, text, max_w, size, fill, bold=False, line_spacing=1.12):
    x, y = xy
    fnt = font(size, bold)
    lines = []
    for block in text.split("\n"):
        if not block:
            lines.append("")
            continue
        current = ""
        for word in block.split():
            probe = f"{current} {word}".strip()
            if draw.textbbox((0, 0), probe, font=fnt)[2] <= max_w:
                current = probe
            else:
                if current:
                    lines.append(current)
                current = word
        if current:
            lines.append(current)
    line_h = int(size * line_spacing)
    for line in lines:
        draw.text((x, y), line, font=fnt, fill=fill)
        y += line_h
    return y


def preview_card(draw, x, y, w, h, title, accent):
    draw.rounded_rectangle([x, y, x + w, y + h], radius=18, fill="white", outline="#ccd6dc", width=2)
    draw.rectangle([x, y, x + w, y + 24], fill=accent)
    draw.text((x + 22, y + 36), title.upper(), font=font(22, True), fill=accent)
    return x + 22, y + 82, w - 44, h - 100


def paste_fit(canvas, path: Path, box):
    x, y, w, h = box
    img = Image.open(path).convert("RGB")
    img.thumbnail((w, h), Image.Resampling.LANCZOS)
    canvas.paste(img, (x + (w - img.width) // 2, y + (h - img.height) // 2))


def build_preview() -> None:
    extract_header()
    scale = 80
    W, H = int(SLIDE_W * scale), int(SLIDE_H * scale)
    img = Image.new("RGB", (W, H), rgb_hex(COLORS["paper"]))
    draw = ImageDraw.Draw(img)

    header = Image.open(HEADER_IMG).convert("RGB").resize((W, int(2.72 * scale)), Image.Resampling.LANCZOS)
    img.paste(header, (0, 0))

    def sx(v): return int(v * scale)

    draw_wrapped(draw, (sx(0.85), sx(3.05)), "Sistema de detección de intrusiones para redes IoT/MQTT", sx(19.8), 48, rgb_hex(COLORS["green_dark"]), True)
    draw_wrapped(draw, (sx(0.85), sx(3.95)), "TinyML en ESP32-S3 + Aprendizaje Federado Jerárquico + ASCON-128 en arquitectura Edge-Fog-Cloud", sx(18.7), 25, rgb_hex(COLORS["ink"]))
    draw_wrapped(draw, (sx(20.4), sx(3.22)), "Santiago Alejandro Jaimes Puerto\nNicolás Casas Ibarra\nDirector: Prof. Carlos Andrés Lozano Garzón, PhD", sx(6.1), 17, rgb_hex(COLORS["muted"]))

    metrics = [("FEATURES", "13", "flujos MQTT/TCP"), ("CLASES", "3", "normal/brute force/scan"), ("MODELO EDGE", "4.5 KB", "MLP/RN en ESP32"), ("FEDAVG", "163", "parámetros federados"), ("RONDA LOCAL", "40", "muestras + 5 épocas"), ("ASCON", "0.06%", "overhead temporal")]
    for i, (label, value, note) in enumerate(metrics):
        x = sx(0.85 + i * (4.05 + 0.28))
        y = sx(5.05)
        w = sx(4.05)
        draw.rounded_rectangle([x, y, x + w, y + sx(1.65)], radius=15, fill="white", outline="#ccd6dc", width=2)
        tw = draw.textbbox((0, 0), label, font=font(13, True))[2]
        draw.text((x + (w - tw) // 2, y + sx(0.2)), label, font=font(13, True), fill=rgb_hex(COLORS["muted"]))
        tw = draw.textbbox((0, 0), value, font=font(36, True))[2]
        draw.text((x + (w - tw) // 2, y + sx(0.52)), value, font=font(36, True), fill=rgb_hex(COLORS["orange"] if i == 5 else COLORS["green_mid"]))
        tw = draw.textbbox((0, 0), note, font=font(12))[2]
        draw.text((x + (w - tw) // 2, y + sx(1.22)), note, font=font(12), fill=rgb_hex(COLORS["muted"]))

    # Preview text blocks mirror the PPTX at lower fidelity.
    col_y = sx(7.25)
    left_x, mid_x, right_x = sx(0.85), sx(9.55), sx(18.25)
    col_w, right_w = sx(8.25), sx(8.45)

    cards = [
        (left_x, col_y, col_w, sx(4.35), "Problema y objetivo", COLORS["green"], "Los IDS centralizados logran alta precisión, pero no son realistas para IoT restringido.\n\nObjetivo: validar un IDS distribuido que detecte ataques MQTT en el borde, entrene sin mover datos crudos y proteja pesos/features."),
        (left_x, sx(12.05), col_w, sx(5.35), "Pipeline experimental", COLORS["blue"], "• Dataset de 256 276 flujos con 13 features.\n• Triclase: normal, mqtt_bruteforce y scan_A.\n• ESP32-S3 infiere y publica features.\n• Raspberry Pi etiqueta, entrena y agrega.\n• PC coordina FedAvg y dashboard."),
        (left_x, sx(17.9), col_w, sx(5.75), "Modelos evaluados", COLORS["green"], "MLP/RN: 13 -> 32 -> 16 -> 8 -> 3; 1 139 parámetros; 4.5 KB.\n\nCNN-1D: Conv1D + BN + GAP + Dense; federación de capas densas finales."),
        (left_x, sx(24.15), col_w, sx(5.45), "Seguridad y variantes", COLORS["orange"], "• ASCON-128: payload {ct, tag, nonce}.\n• No-ASCON: JSON plano con rutas *_plain.\n• FOG: leader/peer antes del PC.\n• Comparación: RN+ASCON, RN sin ASCON, CNN-1D y CNN-1D+FOG."),
        (left_x, sx(30.1), col_w, sx(6.55), "Aporte de la tesis", COLORS["green"], "• Ciclo HFL bidireccional en hardware real.\n• TinyML + MQTT + FastAPI + FedAvg + ASCON.\n• Métricas de convergencia, pesos y overhead.\n• Aísla cifrado, modelo y topología FOG."),
        (mid_x, col_y, col_w, sx(7.85), "Arquitectura implementada", COLORS["green"], "PC Cloud\nmodelo global baja / pesos suben\nRaspberry Pi FOG Leader <-> Raspberry Pi FOG Peer\nMQTT hacia nodos edge\nESP32-S3 Edge A            ESP32-S3 Edge B\n\nASCON: {ct, tag, nonce}. Sin ASCON: JSON plano *_plain."),
        (mid_x, sx(15.6), col_w, sx(6.4), "Flujo bidireccional v7", COLORS["blue"], "1. PC despliega modelo global.\n2. Gateway publica modelo a ESP32.\n3. ESP32 infiere y publica features.\n4. Gateway entrena con 40 muestras.\n5. Sin FOG: reporta al PC.\n6. Con FOG: leader preagrega peers.\n7. PC actualiza el modelo global."),
        (mid_x, sx(22.55), col_w, sx(5.8), "Contraste de rutas", COLORS["orange"], "Con ASCON se usan topics base y payload cifrado. Sin ASCON se usan rutas *_plain para medir el baseline sin costo criptográfico.\n\nRN+ASCON 94.63% | RN sin ASCON 92.98% | CNN FOG 97.50%"),
        (mid_x, sx(28.85), col_w, sx(7.8), "Hallazgos de operación", COLORS["green"], "• HFL converge alrededor de 90-95%.\n• ASCON no degrada la convergencia.\n• FOG permite preagregar gateways.\n• CNN-1D+FOG obtiene mayor accuracy observado.\n• Comunicación: 7.2 KB plain vs 9.8 KB ASCON."),
        (right_x, col_y, right_w, sx(5.2), "Resultados HFL v7", COLORS["green"], "Accuracy global final promedio:\nRN+ASCON: 94.63%\nRN sin ASCON: 92.98%\nCNN-1D+FOG: 97.50%"),
        (right_x, sx(30.15), right_w, sx(6.5), "Conclusión", COLORS["green"], "TinyML + HFL + ASCON-128 es viable para IDS IoT/MQTT. El sistema evita centralizar datos crudos, actualiza modelos de forma bidireccional y agrega seguridad autenticada con costo temporal despreciable."),
    ]
    for x, y, w, h, title, accent, text in cards:
        ix, iy, iw, _ = preview_card(draw, x, y, w, h, title, rgb_hex(accent))
        draw_wrapped(draw, (ix, iy), text, iw, 17 if h < sx(6) else 18, rgb_hex(COLORS["ink"]), False)

    # Result figures.
    ix, iy, iw, ih = preview_card(draw, right_x, sx(12.95), right_w, sx(8.05), "Evolución de accuracy y loss", rgb_hex(COLORS["blue"]))
    paste_fit(img, ACC_LOSS_IMG, (ix, iy, iw, ih - 10))
    ix, iy, iw, ih = preview_card(draw, right_x, sx(21.55), right_w, sx(8.05), "Evolución de magnitud de pesos", rgb_hex(COLORS["orange"]))
    paste_fit(img, WEIGHTS_IMG, (ix, iy, iw, ih - 10))

    draw.rectangle([0, sx(37.85), W, sx(38.85)], fill=rgb_hex(COLORS["green_dark"]))
    draw_wrapped(draw, (sx(0.85), sx(38.08)), "Artefactos: hfl_v7-RN | hfl_v7-no-ascon | hfl_v7-CNN | Documento_Final | Analisis de Modelos", sx(18), 14, "white")
    img.save(OUT_PREVIEW)


if __name__ == "__main__":
    build_pptx()
    build_preview()
    print(f"PPTX: {OUT_PPTX}")
    print(f"Preview: {OUT_PREVIEW}")
