import os
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = Path(os.environ.get("HFL_PPT_OUT", ROOT / "Presentacion_Avances_HFL_v7.pptx"))
PREVIEW_DIR = ROOT / "presentation_previews"
LOGO = ROOT / "HojaTitulo" / "LogoUniandes.png"
MATRIX_ACCURACY_LOSS = ROOT / "images" / "Resultados" / "hfl_v7_accuracy_loss_matrix.png"
MATRIX_WEIGHTS = ROOT / "images" / "Resultados" / "hfl_v7_weight_magnitude_matrix.png"


COLORS = {
    "ink": "151515",
    "muted": "5E6470",
    "paper": "F7F4EE",
    "white": "FFFFFF",
    "yellow": "F2C230",
    "teal": "2A9D8F",
    "blue": "26547C",
    "red": "C44536",
    "green": "4B7F52",
    "line": "D9D4C8",
}


def rgb(hex_color):
    hex_color = COLORS.get(hex_color, hex_color)
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_text(slide, text, x, y, w, h, size=24, bold=False, color="ink", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = box.text_frame.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(COLORS[color])
    r.font.name = "Aptos"
    return box


def add_bullets(slide, items, x, y, w, h, size=22, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        r = p.runs[0]
        r.font.size = Pt(size)
        r.font.color.rgb = rgb(COLORS[color])
        r.font.name = "Aptos"
    return box


def add_slide_title(slide, title, kicker=None):
    if kicker:
        add_text(slide, kicker.upper(), 0.55, 0.25, 3.8, 0.28, size=9, bold=True, color="teal")
    add_text(slide, title, 0.55, 0.48, 8.6, 0.62, size=25, bold=True, color="ink")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.18), Inches(1.25), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(COLORS["yellow"])
    line.line.fill.background()


def add_footer(slide, idx):
    add_text(slide, f"Avances HFL v7 | {idx:02d}", 10.85, 7.08, 1.9, 0.2, size=8, color="muted", align=PP_ALIGN.RIGHT)


def add_logo(slide, x=11.55, y=0.2, w=1.1):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), width=Inches(w))


def add_band(slide, x, y, w, h, color, transparency=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[color])
    shp.fill.transparency = transparency
    shp.line.fill.background()
    return shp


def add_round_box(slide, x, y, w, h, color, text, size=18, text_color="white"):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[color])
    shp.line.color.rgb = rgb(COLORS[color])
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = rgb(COLORS[text_color])
    r.font.name = "Aptos"
    return shp


def add_arrow(slide, x1, y1, x2, y2, color="muted"):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = rgb(COLORS[color])
    conn.line.width = Pt(2.0)
    return conn


def add_table(slide, rows, x, y, w, h, col_widths=None, font_size=11):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for idx, cw in enumerate(col_widths):
            table.columns[idx].width = Inches(cw)
    for r_i, row in enumerate(rows):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = str(val)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r_i == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(font_size)
                    run.font.bold = r_i == 0
                    run.font.color.rgb = rgb(COLORS["white"] if r_i == 0 else COLORS["ink"])
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(COLORS["blue"] if r_i == 0 else ("white" if r_i % 2 else "paper"))
    return table_shape


def slide_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS["paper"])


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 cover
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_band(slide, 0, 0, 13.333, 7.5, "ink")
    add_band(slide, 0, 5.85, 13.333, 1.65, "yellow")
    add_text(slide, "Avances HFL v7", 0.7, 0.75, 4.5, 0.6, size=20, bold=True, color="yellow")
    add_text(slide, "IDS IoT con TinyML,\naprendizaje federado\ny ASCON-128", 0.7, 1.45, 7.0, 2.25, size=38, bold=True, color="white")
    add_text(slide, "ESP32-S3 -> Raspberry Pi -> PC | RN, CNN, no-ASCON y FOG", 0.75, 5.98, 8.6, 0.35, size=17, bold=True, color="ink")
    add_text(slide, "Santiago Alejandro Jaimes Puerto · Nicolas Casas Ibarra", 0.75, 6.45, 6.8, 0.25, size=12, color="ink")
    add_logo(slide, 11.4, 6.08, 1.05)
    for i, (txt, col) in enumerate([("Edge", "teal"), ("Fog", "yellow"), ("Cloud", "red")]):
        add_round_box(slide, 8.4 + i * 1.35, 2.65 + (i % 2) * 0.45, 1.05, 0.52, col, txt, size=14, text_color="ink" if col == "yellow" else "white")

    # 2 delta from v5
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Del plan v5 al sistema v7 implementado", "Cambio de alcance")
    add_logo(slide)
    add_text(slide, "Antes", 0.8, 1.65, 1.2, 0.3, 18, True, "muted")
    add_text(slide, "Ahora", 7.1, 1.65, 1.2, 0.3, 18, True, "teal")
    add_round_box(slide, 0.9, 2.25, 1.7, 0.55, "blue", "PC", 15)
    add_round_box(slide, 2.95, 2.25, 1.7, 0.55, "blue", "RPi", 15)
    add_round_box(slide, 1.9, 3.25, 1.8, 0.55, "blue", "ESP32 broker", 13)
    add_round_box(slide, 3.95, 3.25, 1.8, 0.55, "blue", "ESP32 client", 13)
    add_arrow(slide, 2.6, 2.53, 2.95, 2.53)
    add_arrow(slide, 3.75, 2.85, 3.0, 3.25)
    add_arrow(slide, 3.7, 3.53, 3.95, 3.53)
    add_round_box(slide, 6.8, 2.15, 1.6, 0.55, "teal", "ESP32-S3", 13)
    add_round_box(slide, 8.75, 2.15, 1.9, 0.55, "yellow", "Raspberry Pi", 13, "ink")
    add_round_box(slide, 11.0, 2.15, 1.25, 0.55, "red", "PC", 13)
    add_arrow(slide, 8.4, 2.43, 8.75, 2.43)
    add_arrow(slide, 10.65, 2.43, 11.0, 2.43)
    add_round_box(slide, 8.65, 3.15, 1.25, 0.48, "blue", "FOG", 12)
    add_bullets(slide, [
        "10 clases y cliente separado quedaron como antecedente.",
        "v7 trabaja con 3 clases MQTT y 13 features de flujo.",
        "La comparación real cambia modelo, cifrado y topología."
    ], 0.85, 4.55, 5.3, 1.35, size=17)
    add_bullets(slide, [
        "RN/MLP con ASCON como rama principal.",
        "Baseline sin ASCON para medir overhead.",
        "CNN-1D y FOG como variantes de arquitectura."
    ], 6.85, 4.55, 5.3, 1.35, size=17)
    add_footer(slide, 2)

    # 3 architecture
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Arquitectura HFL bidireccional", "Edge--Fog--Cloud")
    add_logo(slide)
    xs = [0.85, 4.75, 8.85]
    labels = [("Edge", "ESP32-S3\nInferencia TinyML\nMQTT features"), ("Fog", "Raspberry Pi\nEntrenamiento local\nBroker + FedAvg fog"), ("Cloud", "PC FastAPI\nFedAvg global\nDashboard + CSV")]
    for x, (head, body) in zip(xs, labels):
        add_round_box(slide, x, 2.0, 2.55, 0.65, "blue" if head == "Cloud" else ("teal" if head == "Edge" else "yellow"), head, 18, "white" if head != "Fog" else "ink")
        add_text(slide, body, x + 0.1, 2.85, 2.35, 1.05, 16, False, "ink", PP_ALIGN.CENTER)
    add_arrow(slide, 3.45, 2.35, 4.75, 2.35, "teal")
    add_arrow(slide, 7.35, 2.35, 8.85, 2.35, "teal")
    add_arrow(slide, 8.85, 4.3, 7.35, 4.3, "red")
    add_arrow(slide, 4.75, 4.3, 3.45, 4.3, "red")
    add_text(slide, "Bottom-up: features -> entrenamiento local -> pesos", 2.0, 1.55, 5.4, 0.25, 14, True, "teal")
    add_text(slide, "Top-down: modelo global -> gateway -> ESP32", 6.05, 4.58, 5.2, 0.25, 14, True, "red")
    add_bullets(slide, [
        "Topics: fl/features, fl/global_model",
        "Endpoints: /aggregate-from-gateway, /deploy-model",
        "FOG: fog/weights y fog/global_model"
    ], 1.0, 5.45, 11.0, 0.9, size=18)
    add_footer(slide, 3)

    # 4 data
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Datos y problema de clasificación", "Dataset triclase")
    add_logo(slide)
    rows = [
        ["Clase", "Archivo", "Flujos", "Señal principal"],
        ["normal", "uniflow_normal.csv", "171,837", "MQTT/TCP benigno"],
        ["mqtt_bruteforce", "uniflow_mqtt_bruteforce.csv", "33,080", "PSH y paquetes altos"],
        ["scan_A", "uniflow_scan_A.csv", "51,359", "Pocos paquetes, tamaño bajo"],
    ]
    add_table(slide, rows, 0.75, 1.55, 6.7, 1.85, [1.45, 2.5, 1.1, 1.65], 10)
    add_text(slide, "13 features", 8.4, 1.5, 2.0, 0.5, 31, True, "teal")
    add_bullets(slide, [
        "num_pkts, num_bytes",
        "mean/std/min/max IAT",
        "mean/std/min/max packet length",
        "PSH, RST y URG flags"
    ], 8.15, 2.25, 3.7, 1.3, size=17)
    add_text(slide, "Etiquetado online", 0.8, 4.35, 2.2, 0.3, 18, True, "ink")
    add_bullets(slide, [
        "pkts >= 50 y PSH >= 10 -> bruteforce",
        "pkts <= 5, pkt_len <= 50 y PSH <= 1 -> scan_A",
        "pkts <= 30 y pkt_len >= 50 -> normal"
    ], 0.95, 4.85, 10.9, 0.9, size=17)
    add_footer(slide, 4)

    # 5 model comparison
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Modelos probados antes del despliegue", "Notebook general")
    add_logo(slide)
    chart_data = CategoryChartData()
    chart_data.categories = ["MLP/RN", "Residual", "CNN-1D", "Transformer"]
    chart_data.add_series("Accuracy", [0.9060, 0.9055, 0.9051, 0.8903])
    chart_data.add_series("F1 weighted", [0.9049, 0.9044, 0.9041, 0.8919])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.75), Inches(1.55), Inches(7.0), Inches(4.4), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.minimum_scale = 0.86
    chart.value_axis.maximum_scale = 0.92
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Desempeño offline"
    for series, color in zip(chart.series, ["blue", "red"]):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = rgb(COLORS[color])
    add_bullets(slide, [
        "MLP/RN queda como rama principal por balance.",
        "CNN-1D se implementa para comparar arquitectura.",
        "Transformer se conserva como experimento offline.",
        "Todos exportan pesos, scaler y label_map."
    ], 8.15, 1.75, 4.1, 2.2, size=18)
    add_text(slide, "Decisión: federar solo las capas finales para reducir comunicación.", 8.15, 5.25, 4.0, 0.65, 19, True, "blue")
    add_footer(slide, 5)

    # 6 federated process
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Entrenamiento federado implementado", "Proceso por ronda")
    add_logo(slide)
    steps = [
        ("1", "ESP32 envía features cada 5 s"),
        ("2", "Gateway acumula 40 muestras"),
        ("3", "Entrena 5 épocas, batch 8"),
        ("4", "Envía 163 parámetros"),
        ("5", "PC aplica FedAvg"),
        ("6", "Modelo global vuelve al edge"),
    ]
    for i, (num, txt) in enumerate(steps):
        x = 0.75 + (i % 3) * 4.05
        y = 1.65 + (i // 3) * 2.0
        add_round_box(slide, x, y, 0.48, 0.48, "yellow", num, 14, "ink")
        add_text(slide, txt, x + 0.65, y - 0.02, 2.75, 0.55, 17, True, "ink")
        if i not in (2, 5):
            add_arrow(slide, x + 3.45, y + 0.25, x + 3.9, y + 0.25, "muted")
    add_text(slide, "El modelo completo tiene 1,139 parámetros; el ciclo federado comparte solo W3, b3, W4 y b4.", 0.95, 6.0, 10.8, 0.45, 19, True, "blue", PP_ALIGN.CENTER)
    add_footer(slide, 6)

    # 7 ASCON
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "ASCON-128 protege el ciclo completo", "Seguridad")
    add_logo(slide)
    rows = [
        ["Canal", "Contenido", "Modo RN", "Baseline"],
        ["ESP32 -> RPi", "Features", "MQTT + ASCON", "MQTT plain"],
        ["RPi -> PC", "Pesos + métricas", "HTTP + ASCON", "HTTP plain"],
        ["PC -> RPi", "Modelo global", "HTTP + ASCON", "HTTP plain"],
        ["RPi -> ESP32", "Modelo global", "MQTT + ASCON", "MQTT plain"],
    ]
    add_table(slide, rows, 0.65, 1.45, 7.45, 2.55, [1.45, 1.9, 2.0, 1.6], 9)
    add_text(slide, "< 0.1%", 8.85, 1.55, 2.1, 0.55, 36, True, "teal")
    add_text(slide, "overhead temporal relativo por ronda", 8.9, 2.18, 3.0, 0.35, 15, True, "ink")
    add_text(slide, "~36%", 8.85, 3.25, 2.1, 0.55, 36, True, "red")
    add_text(slide, "incremento de tamaño por mensaje", 8.9, 3.88, 3.0, 0.35, 15, True, "ink")
    add_text(slide, "La rama no-ASCON permite medir el costo del cifrado sin cambiar el resto de la arquitectura.", 1.0, 5.45, 10.7, 0.6, 18, True, "blue", PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8 fog
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Estrategia FOG", "Agregación intermedia")
    add_logo(slide)
    add_round_box(slide, 0.9, 2.4, 1.8, 0.6, "teal", "Gateway A\nleader", 13)
    add_round_box(slide, 0.9, 4.45, 1.8, 0.6, "teal", "Gateway B\npeer", 13)
    add_round_box(slide, 5.3, 3.35, 2.0, 0.75, "yellow", "Fog FedAvg", 18, "ink")
    add_round_box(slide, 10.0, 3.35, 1.8, 0.75, "blue", "PC\nCloud", 16)
    add_arrow(slide, 2.75, 2.7, 5.3, 3.55, "teal")
    add_arrow(slide, 2.75, 4.75, 5.3, 3.85, "teal")
    add_arrow(slide, 7.3, 3.72, 10.0, 3.72, "teal")
    add_arrow(slide, 10.0, 4.25, 7.3, 4.25, "red")
    add_text(slide, "fog/weights", 3.4, 2.65, 1.4, 0.25, 12, True, "teal")
    add_text(slide, "/aggregate-from-fog", 7.8, 3.35, 1.8, 0.25, 12, True, "teal")
    add_text(slide, "fog/global_model", 7.65, 4.42, 1.8, 0.25, 12, True, "red")
    add_bullets(slide, [
        "El líder agrega pesos locales y pesos peer.",
        "El PC recibe una actualización preagregada.",
        "El modelo global se redistribuye a peers y ESP32."
    ], 0.95, 5.8, 10.9, 0.7, size=17)
    add_footer(slide, 8)

    # 9 variant results
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Resultados de las variantes v7", "Evidencia experimental")
    add_logo(slide)
    chart_data = CategoryChartData()
    chart_data.categories = ["RN + ASCON", "RN sin ASCON", "CNN FOG"]
    chart_data.add_series("Accuracy final", [0.9463, 0.9298, 0.9750])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.75), Inches(1.45), Inches(6.2), Inches(4.2), chart_data).chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0.88
    chart.value_axis.maximum_scale = 1.0
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Accuracy global final promedio"
    for point, color in zip(chart.series[0].points, ["blue", "red", "green"]):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = rgb(COLORS[color])
    rows = [
        ["Variante", "Intentos", "Loss", "Ronda"],
        ["RN + ASCON", "9", "0.1415", "61.25 s"],
        ["RN sin ASCON", "14", "0.1581", "51.01 s"],
        ["CNN FOG", "7", "0.0910", "58.61 s"],
    ]
    add_table(slide, rows, 7.3, 1.65, 4.9, 2.25, [1.8, 0.8, 0.9, 1.1], 10)
    add_text(slide, "Lectura: ASCON no rompe la convergencia; FOG y CNN funcionan, pero el MLP/RN conserva la ruta más simple para defender.", 7.35, 4.55, 4.55, 0.9, 17, True, "blue")
    add_footer(slide, 9)

    # 10 accuracy/loss curves
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Curvas de accuracy y loss", "Comparación temporal 2x2")
    add_logo(slide)
    slide.shapes.add_picture(str(MATRIX_ACCURACY_LOSS), Inches(0.55), Inches(1.28), height=Inches(5.7))
    add_bullets(slide, [
        "RN + ASCON y no-ASCON mantienen accuracy alto con oscilaciones controladas.",
        "CNN y CNN_FOG muestran recuperación progresiva después de las primeras rondas.",
        "La lectura temporal evita depender solo del último punto de cada corrida."
    ], 9.0, 1.7, 3.6, 2.25, size=17)
    add_text(slide, "La matriz compara dinámica, no solo promedio final.", 9.05, 5.45, 3.35, 0.5, 18, True, "blue")
    add_footer(slide, 10)

    # 11 weight curves
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Magnitud de pesos", "Estabilidad de actualizaciones")
    add_logo(slide)
    slide.shapes.add_picture(str(MATRIX_WEIGHTS), Inches(0.55), Inches(1.28), height=Inches(5.7))
    add_bullets(slide, [
        "Las curvas se mantienen acotadas en las cuatro variantes.",
        "CNN/CNN_FOG desplazan más la magnitud de pesos de salida.",
        "No se observa divergencia durante las rondas analizadas."
    ], 9.0, 1.7, 3.6, 2.25, size=17)
    add_text(slide, "Sirve como evidencia de estabilidad del ciclo HFL.", 9.05, 5.45, 3.35, 0.5, 18, True, "teal")
    add_footer(slide, 11)

    # 12 paper updates
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Paper actualizado", "Cambios incorporados")
    add_logo(slide)
    add_bullets(slide, [
        "Introducción y objetivos alineados con v7 real: ESP32-S3, Raspberry Pi 4 y PC.",
        "Marco teórico depurado: TinyML, HFL, ASCON y FOG.",
        "Metodología ampliada: RN/MLP, CNN, no-ASCON, FOG y notebook de entrenamiento.",
        "Resultados actualizados con comparación offline y métricas de variantes.",
        "Conclusiones ajustadas: arquitectura modular y validación del ciclo bidireccional."
    ], 0.95, 1.65, 11.0, 3.2, size=20)
    add_text(slide, "El PDF ya compila con Tectonic: build/EncabezadoTesisMSc.pdf", 0.95, 5.85, 10.8, 0.35, 18, True, "teal")
    add_footer(slide, 12)

    # 13 next
    slide = prs.slides.add_slide(blank)
    slide_background(slide)
    add_slide_title(slide, "Mensaje para la reunión", "Síntesis")
    add_logo(slide)
    add_text(slide, "Lo que ya podemos mostrar", 0.95, 1.55, 4.5, 0.45, 24, True, "ink")
    add_bullets(slide, [
        "HFL bidireccional funcionando en Edge--Fog--Cloud.",
        "Tres variantes comparables: RN, no-ASCON y CNN.",
        "Modo FOG implementado para agregación entre gateways.",
        "ASCON medido contra baseline y viable en tiempo."
    ], 1.0, 2.2, 5.3, 2.0, size=20)
    add_text(slide, "Pendientes honestos", 7.05, 1.55, 4.5, 0.45, 24, True, "red")
    add_bullets(slide, [
        "Mejorar redacción final y uniformar nombres de ramas.",
        "Reducir overfull boxes del LaTeX si se entrega versión final.",
        "Agregar fotos/hardware real si el profesor pide evidencia visual.",
        "Cerrar pendientes de redacción, tablas y anexos."
    ], 7.1, 2.2, 5.4, 2.0, size=20)
    add_band(slide, 0.0, 6.35, 13.333, 1.15, "ink")
    add_text(slide, "Tesis técnica: un IDS IoT federado, cifrado y modular puede operar sobre hardware de borde real.", 0.9, 6.58, 11.6, 0.4, 20, True, "white", PP_ALIGN.CENTER)
    add_footer(slide, 13)

    prs.save(OUT)
    return prs


def load_font(size, bold=False):
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/calibrib.ttf" if bold else "C:/Windows/Fonts/calibri.ttf",
    ]
    for c in candidates:
        if Path(c).exists():
            return ImageFont.truetype(c, size)
    return ImageFont.load_default()


def wrap_text(draw, text, font, max_width):
    words = text.split()
    lines = []
    line = ""
    for word in words:
        test = f"{line} {word}".strip()
        if draw.textbbox((0, 0), test, font=font)[2] <= max_width:
            line = test
        else:
            if line:
                lines.append(line)
            line = word
    if line:
        lines.append(line)
    return lines


def generate_previews():
    PREVIEW_DIR.mkdir(exist_ok=True)
    size = (1600, 900)
    slide_titles = [
        ("Avances HFL v7", "IDS IoT con TinyML, HFL y ASCON-128"),
        ("Del plan v5 al sistema v7", "3 capas operativas + FOG opcional"),
        ("Arquitectura HFL bidireccional", "Edge -> Fog -> Cloud -> Fog -> Edge"),
        ("Datos y problema", "3 clases, 13 features, 256k flujos"),
        ("Modelos probados", "MLP/RN, CNN-1D, Residual-MLP, Transformer"),
        ("Entrenamiento federado", "40 muestras, 5 épocas, 163 parámetros"),
        ("ASCON-128", "Seguridad con overhead temporal bajo"),
        ("Estrategia FOG", "Leader/peer y FedAvg intermedio"),
        ("Resultados v7", "RN + ASCON, no-ASCON, CNN FOG"),
        ("Curvas accuracy/loss", "Matriz 2x2 por variante"),
        ("Magnitud de pesos", "Matriz 2x2 de estabilidad"),
        ("Paper actualizado", "Objetivos, metodología, resultados y conclusiones"),
        ("Mensaje para la reunión", "Arquitectura validada y pendientes honestos"),
    ]
    title_font = load_font(48, True)
    sub_font = load_font(27, False)
    small_font = load_font(18, False)
    for idx, (title, subtitle) in enumerate(slide_titles, 1):
        img = Image.new("RGB", size, "#" + COLORS["paper"])
        d = ImageDraw.Draw(img)
        d.rectangle([0, 0, 1600, 96], fill="#" + COLORS["ink"])
        d.rectangle([0, 830, 1600, 900], fill="#" + (COLORS["yellow"] if idx == 1 else COLORS["line"]))
        d.text((70, 30), f"{idx:02d}", font=small_font, fill="#" + COLORS["yellow"])
        d.text((140, 150), title, font=title_font, fill="#" + COLORS["ink"])
        for j, line in enumerate(wrap_text(d, subtitle, sub_font, 1050)):
            d.text((140, 230 + j * 38), line, font=sub_font, fill="#" + COLORS["muted"])
        if idx in (3, 6, 8):
            for k, label in enumerate(["Edge", "Fog", "Cloud"]):
                x = 180 + k * 430
                d.rounded_rectangle([x, 430, x + 250, 505], radius=18, fill="#" + [COLORS["teal"], COLORS["yellow"], COLORS["blue"]][k])
                d.text((x + 72, 452), label, font=sub_font, fill="#" + ("151515" if label == "Fog" else "FFFFFF"))
                if k < 2:
                    d.line([x + 255, 467, x + 390, 467], fill="#" + COLORS["muted"], width=5)
        elif idx in (5, 9, 10, 11):
            vals = [0.906, 0.905, 0.905, 0.890] if idx == 5 else ([0.946, 0.930, 0.975] if idx == 9 else [0.93, 0.95, 0.91, 0.97])
            labels = ["MLP", "Residual", "CNN", "Transf"] if idx == 5 else (["RN", "Plain", "CNN FOG"] if idx == 9 else ["RN", "Plain", "CNN", "FOG"])
            for k, v in enumerate(vals):
                x = 180 + k * 260
                h = int((v - 0.86) / 0.14 * 280)
                d.rectangle([x, 710 - h, x + 150, 710], fill="#" + [COLORS["teal"], COLORS["green"], COLORS["blue"], COLORS["red"]][k % 4])
                d.text((x, 725), labels[k], font=small_font, fill="#" + COLORS["ink"])
                d.text((x, 675 - h), f"{v:.3f}", font=small_font, fill="#" + COLORS["ink"])
        else:
            y = 390
            for bullet in ["Arquitectura modular", "Comparación con baseline", "Resultados listos para discusión"]:
                d.ellipse([155, y + 7, 175, y + 27], fill="#" + COLORS["teal"])
                d.text((195, y), bullet, font=sub_font, fill="#" + COLORS["ink"])
                y += 58
        img.save(PREVIEW_DIR / f"slide_{idx:02d}.png")

    montage = Image.new("RGB", (1600, 1320), "#" + COLORS["white"])
    thumbs = []
    for idx in range(1, len(slide_titles) + 1):
        im = Image.open(PREVIEW_DIR / f"slide_{idx:02d}.png").resize((400, 225))
        thumbs.append(im)
    for idx, im in enumerate(thumbs):
        x = (idx % 4) * 400
        y = (idx // 4) * 330
        montage.paste(im, (x, y))
    montage.save(PREVIEW_DIR / "montage.png")


if __name__ == "__main__":
    build_deck()
    generate_previews()
    print(OUT)
    print(PREVIEW_DIR / "montage.png")
